#!/usr/bin/env python3
"""Generate the Computer Vision Project presentation. 17 slides.

Design system
-------------
Palette : navy primary, teal accent, amber signal, indigo (Explain), slate body.
Grid    : single left margin, shared title/caption baselines, consistent rhythm.
Type    : strict 4-level scale (title 30 / subtitle 18 / body 16 / caption 13).
Polish  : white cards with colored top-borders, play-icon video placeholders,
          left-edge section spine, unified teal arrows, generous whitespace.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Canvas ────────────────────────────────────────────────────────────────────
PRS_WIDTH = Inches(13.333)
PRS_HEIGHT = Inches(7.5)

# ── Color system ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x29, 0x3B)   # headings, key text
TEAL        = RGBColor(0x0D, 0x94, 0x88)   # primary accent (Detect)
AMBER       = RGBColor(0xD9, 0x77, 0x06)   # warm signal (Enrich)
INDIGO      = RGBColor(0x63, 0x66, 0xF1)   # Explain section
SLATE       = RGBColor(0x47, 0x55, 0x69)   # body text
LIGHT_SLATE = RGBColor(0x94, 0xA3, 0xB8)   # captions / footnotes
BG          = RGBColor(0xF8, 0xFA, 0xFC)   # slide background
CARD        = RGBColor(0xFF, 0xFF, 0xFF)   # card surface
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)   # hairline dividers
PLACEHOLDER = RGBColor(0xEE, 0xF2, 0xF6)   # video / image placeholder fill
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)   # divider background (slate-900)
DARK_FADE   = RGBColor(0x1E, 0x29, 0x3B)   # giant faded number on dark
DARK_LABEL  = RGBColor(0x94, 0xA3, 0xB8)   # progress label on dark

# ── Grid constants ────────────────────────────────────────────────────────────
MARGIN    = Inches(0.9)     # single left anchor for all content
CONTENT_W = Inches(11.533)  # standard block width (slide - 2*margin)
TITLE_Y   = Inches(0.55)    # shared title baseline
CAPTION_Y = Inches(6.55)    # shared caption baseline
SPINE_W   = Inches(0.14)    # left-edge section bar width

# ── Typography scale ──────────────────────────────────────────────────────────
SZ_TITLE    = Pt(30)
SZ_SUBTITLE = Pt(18)
SZ_BODY     = Pt(16)
SZ_CAPTION  = Pt(13)
FONT        = "Calibri"

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size, font_color,
                bold=False, align=PP_ALIGN.LEFT, anchor=None, font_name=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = font_name
    return tf

def add_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def no_line(shape):
    shape.line.fill.background()

def thin_border(shape, color=BORDER, width=Pt(1)):
    shape.line.color.rgb = color
    shape.line.width = width

def add_title(slide, text, color=NAVY):
    """Slide title on the shared baseline, left-anchored to the grid."""
    return add_textbox(slide, MARGIN, TITLE_Y, CONTENT_W, Inches(0.8),
                       text, SZ_TITLE, color, bold=True)

def add_subtitle(slide, text, top, color=SLATE):
    return add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.5),
                       text, SZ_SUBTITLE, color)

def add_caption(slide, text, color=LIGHT_SLATE, top=CAPTION_Y):
    return add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.5),
                       text, SZ_CAPTION, color)

def section_spine(slide, color):
    """Full-height left-edge color bar that gives each section a visual spine."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 SPINE_W, PRS_HEIGHT)
    bar.fill.solid(); bar.fill.fore_color.rgb = color; no_line(bar)
    return bar

def accent_rule(slide, left, top, width, color=TEAL, height=Pt(3)):
    """A short, deliberate accent line — the single 'pop' per slide."""
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rule.fill.solid(); rule.fill.fore_color.rgb = color; no_line(rule)
    return rule

def video_placeholder(slide, left, top, width, height, label):
    """White card + hairline border + centered play icon + label below."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = PLACEHOLDER
    thin_border(card, BORDER, Pt(1))
    # Play icon: circle + triangle, centered
    d = Inches(0.9)
    cx = left + (width - d) / 2
    cy = top + (height - d) / 2 - Inches(0.25)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    circle.fill.solid(); circle.fill.fore_color.rgb = WHITE
    thin_border(circle, BORDER, Pt(1.5))
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 cx + Inches(0.32), cy + Inches(0.27),
                                 Inches(0.32), Inches(0.36))
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = TEAL; no_line(tri)
    # Label below the icon
    add_textbox(slide, left, cy + d + Inches(0.1), width, Inches(0.8),
                label, SZ_CAPTION, SLATE, align=PP_ALIGN.CENTER)
    return card

def stat_card(slide, left, top, width, height, value, label, accent=TEAL):
    """White card with a colored top-border; the number carries the color."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    thin_border(card, BORDER, Pt(1))
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = accent; no_line(top_bar)
    add_textbox(slide, left, top + Inches(0.18), width, Inches(0.5),
                value, Pt(22), accent, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.72), width, Inches(0.4),
                label, SZ_CAPTION, SLATE, align=PP_ALIGN.CENTER)
    return card

def pill(slide, left, top, width, height, text, color):
    """Small rounded badge for stage/flow labels."""
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    p.fill.solid(); p.fill.fore_color.rgb = WHITE; thin_border(p, color, Pt(1.25))
    add_textbox(slide, left, top + Inches(0.04), width, height,
                text, Pt(11), color, bold=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    return p

def teal_arrow(slide, left, top, width, height, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid(); a.fill.fore_color.rgb = color; no_line(a)
    return a

def divider_slide(prs, blank, word, idx, color):
    """Dark section divider: a deliberate pause between the light content slides.

    Dark navy canvas + colored left spine + giant faded number + white word +
    colored accent rule + progress label. Offers visual breathing room.
    """
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, DARK_BG)
    section_spine(s, color)
    # Giant faded number — background layer, drawn first
    add_textbox(s, MARGIN, Inches(1.05), CONTENT_W, Inches(3.2),
                f"{idx:02d}", Pt(200), DARK_FADE, bold=True, align=PP_ALIGN.CENTER)
    # Small section eyebrow above the word
    add_textbox(s, MARGIN, Inches(1.95), CONTENT_W, Inches(0.4),
                f"SECTION {idx}", Pt(15), color, bold=True, align=PP_ALIGN.CENTER)
    # Section word
    add_textbox(s, MARGIN, Inches(2.7), CONTENT_W, Inches(1.4),
                word, Pt(66), WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Accent rule centered under the word
    accent_rule(s, Inches(5.97), Inches(4.3), Inches(1.4), color)
    # Progress label
    add_textbox(s, MARGIN, Inches(4.55), CONTENT_W, Inches(0.5),
                f"{idx} of 3", Pt(15), DARK_LABEL, align=PP_ALIGN.CENTER)
    return s

# ── Build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = PRS_WIDTH
prs.slide_height = PRS_HEIGHT
blank = prs.slide_layouts[6]

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_textbox(s, MARGIN, Inches(2.35), CONTENT_W, Inches(1.2),
            "Street-Scene Intelligence", Pt(52), NAVY, bold=True, align=PP_ALIGN.CENTER)
accent_rule(s, Inches(5.67), Inches(3.55), Inches(2.0), TEAL)
add_textbox(s, MARGIN, Inches(3.85), CONTENT_W, Inches(0.9),
            "A Multi-Stage Pipeline for Traffic-Scene Detection, Brand Recognition, and Natural-Language Explanation",
            SZ_SUBTITLE, SLATE, align=PP_ALIGN.CENTER)
add_textbox(s, MARGIN, Inches(5.05), CONTENT_W, Inches(0.5),
            "Dongyuan Gao   ·   Solène Cosandey", Pt(15), LIGHT_SLATE, align=PP_ALIGN.CENTER)
add_note(s, "This is our computer-vision project: a multi-stage pipeline that turns raw dashcam video into structured, narrated, and queryable street understanding.")

# ── SLIDE 2: The Hook ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_textbox(s, MARGIN, Inches(1.9), CONTENT_W, Inches(0.5),
            "Our demo pipeline bridges the gap between raw detection and human trust.",
            SZ_SUBTITLE, SLATE, align=PP_ALIGN.CENTER)
add_textbox(s, MARGIN, Inches(2.6), CONTENT_W, Inches(1.3),
            "How do we improve trust\nin self-driving cars?", Pt(46), NAVY,
            bold=True, align=PP_ALIGN.CENTER)
accent_rule(s, Inches(5.97), Inches(4.35), Inches(1.4), TEAL)
add_textbox(s, MARGIN, Inches(4.65), CONTENT_W, Inches(0.5),
            "In the next 15 minutes:  Detect  →  Enrich  →  Explain", SZ_SUBTITLE,
            TEAL, bold=True, align=PP_ALIGN.CENTER)
add_note(s, "Attention: Self-driving cars detect every obstacle — but passengers still don't know why the car braked.\nBenefit: One pipeline turns raw dashcam video into a marketing-ready story — boxed objects, named brands, plain-language captions — so any viewer can follow what the car sees.\nCredibility: 4 stages across 3 sections · YOLO26s · CLIP · VLM · SegFormer — all running locally.\nDestination: In the next 15 minutes: Detect → Enrich → Explain.")

# ── SLIDE 3: TOC ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title(s, "Three Sections")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
sects = [("1", "Detect", "Find and segment objects", TEAL),
         ("2", "Enrich", "Add fine-grained identity", AMBER),
         ("3", "Explain", "Turn everything into a story", INDIGO)]
card_w = Inches(3.5); gap = Inches(0.52)
total = 3 * card_w + 2 * gap
start_x = (PRS_WIDTH - total) / 2
for i, (num, title, desc, color) in enumerate(sects):
    x = start_x + i * (card_w + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), card_w, Inches(3.1))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; thin_border(card, BORDER, Pt(1))
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.4), card_w, Pt(5))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; no_line(top_bar)
    add_textbox(s, x, Inches(2.7), card_w, Inches(0.9), num, Pt(46), color,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(3.65), card_w, Inches(0.6), title, Pt(26), NAVY,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.25), Inches(4.35), card_w - Inches(0.5), Inches(0.8),
                desc, SZ_BODY, SLATE, align=PP_ALIGN.CENTER)
add_note(s, "The talk follows the actual data flow of the pipeline: first we find and segment objects, then we add fine-grained identity, then we turn everything into a story any viewer can follow.")

# ── SECTION 1 SPINE COLOR = TEAL (Detect) ──────────────────────────────────────
# ── SLIDE 4: The Gap ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, TEAL)
add_title(s, "Finding objects is step one — we also need to see the whole scene.")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
gap_w = Inches(5.1)
video_placeholder(s, MARGIN, Inches(1.85), gap_w, Inches(4.3), "Raw dashcam frame")
video_placeholder(s, Inches(7.28), Inches(1.85), gap_w, Inches(4.3), "YOLO bounding boxes")
teal_arrow(s, Inches(6.18), Inches(3.85), Inches(0.95), Inches(0.45))
add_note(s, "YOLO gives us boxes, but a car has a brand. The scene is not just objects — it is road, sidewalk, sky. We need two detection layers.")

# ── SLIDE 5: YOLO Detection ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, TEAL)
add_title(s, "YOLO26s — fast object detection")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
add_subtitle(s, "11 classes · 512 px · 30 epochs · mAP@0.5 = 0.842  ·  mAP@0.5:0.95 = 0.515", Inches(1.55))
video_placeholder(s, MARGIN, Inches(2.2), Inches(7.8), Inches(4.0),
                  "Stage 1 — YOLO-only output (5-sec clip, no brand chips yet)")
for i, (label, val) in enumerate([("Classes", "11"), ("Image Size", "512 px"), ("Epochs", "30")]):
    stat_card(s, Inches(9.05), Inches(2.2) + i * Inches(1.35), Inches(3.35), Inches(1.15),
              val, label, TEAL)
add_caption(s, "Upgraded from YOLOv10n → YOLO26s:  +24.7% mAP@0.5,  +32.7% mAP@0.5:0.95,  +26.1% recall  (11-class Self-Driving-Car-3, 30 epochs).")
add_note(s, "YOLOv10n baseline: P=0.801, R=0.598, mAP@0.5=0.675, mAP@0.5:0.95=0.388.  YOLO26s: P=0.874, R=0.754, mAP@0.5=0.842, mAP@0.5:0.95=0.515.  The 9.4 M parameter backbone with C2PSA blocks delivers a significant accuracy lift — especially on recall and strict mAP. Pedestrian remains the hardest class (mAP 0.692) due to scale and occlusion. The dataset covers 11 street classes including nuanced traffic-light states.")

# ── SLIDE 6: Segmentation ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, TEAL)
add_title(s, "SegFormer-B5 — pixel-level scene understanding")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
add_subtitle(s, "Cityscapes · 19 classes · 84 % mIoU", Inches(1.55))
# Real SegFormer outputs: original, mask, blended overlay
seg_imgs = [
    ("/home/dongyuan/Desktop/computer_vision/other_document/segformer_original.png", "Original"),
    ("/home/dongyuan/Desktop/computer_vision/other_document/segformer_mask.png", "Mask"),
    ("/home/dongyuan/Desktop/computer_vision/other_document/segformer_blended.png", "Overlay (α=0.5)")
]
seg_w = Inches(3.7); seg_gap = Inches(0.22)
for i, (img_path, label) in enumerate(seg_imgs):
    x = MARGIN + i * (seg_w + seg_gap)
    pic = s.shapes.add_picture(img_path, x, Inches(2.2), width=seg_w)
    thin_border(pic, BORDER, Pt(1))
    add_textbox(s, x, Inches(6.15), seg_w, Inches(0.5),
                label, Pt(12), TEAL, bold=True, align=PP_ALIGN.CENTER)
add_caption(s, "Nearest-neighbor upsampling preserves boundaries; 2,516 frames at ~5 it/s (~8 min) — quality demo, not real-time.")
add_note(s, "Our newest stage answers the same question as YOLO, but at the pixel level: what is in this scene? Road (purple), vegetation (green), vehicles (red), sky (blue) — scene geometry, not just boxes. SegFormer-B5 is pretrained on Cityscapes (84% mIoU) and generalises directly to our dashcam domain without any fine-tuning.")

# ── SLIDE 7: Divider Enrich ───────────────────────────────────────────────────
s = divider_slide(prs, blank, "Enrich", 2, AMBER)
add_note(s, "So we can detect cars and map the scene. But knowing that it's a car isn't enough — what kind of car is it?")

# ── SECTION 2 SPINE COLOR = AMBER (Enrich) ──────────────────────────────────────
# ── SLIDE 8: From Generic to Specific ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, AMBER)
add_title(s, "YOLO says 'car.'  We want to say 'Volkswagen.'")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), AMBER)
b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(2.6), Inches(2.4), Inches(1.8))
b1.fill.solid(); b1.fill.fore_color.rgb = PLACEHOLDER; thin_border(b1, TEAL, Pt(1.5))
add_textbox(s, Inches(1.6), Inches(2.6), Inches(2.4), Inches(1.8), "YOLO box",
            SZ_BODY, SLATE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
teal_arrow(s, Inches(4.35), Inches(3.25), Inches(1.5), Inches(0.55), AMBER)
b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(2.1), Inches(4.4), Inches(2.8))
b2.fill.solid(); b2.fill.fore_color.rgb = PLACEHOLDER; thin_border(b2, AMBER, Pt(1.5))
add_textbox(s, Inches(6.3), Inches(2.1), Inches(4.4), Inches(2.8), "High-res car crop",
            SZ_BODY, SLATE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_caption(s, "Use case: fleet analytics, insurance claims, recall targeting, customer trust.", top=Inches(5.3))
add_note(s, "Detection gives us generic classes. Enrichment gives us identity — turning 'car' into 'Toyota' so the viewer knows the system truly understands, not just locates.")

# ── SLIDE 9: CLIP Linear Probe ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, AMBER)
add_title(s, "Frozen backbone + a learned linear probe")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), AMBER)
add_subtitle(s, "OpenCLIP ViT-B-32 (frozen) + nn.Linear(512 → 20) = 20-brand classifier", Inches(1.55))
# Backbone card
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.5), Inches(3.0), Inches(2.4))
bb.fill.solid(); bb.fill.fore_color.rgb = CARD; thin_border(bb, BORDER, Pt(1))
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.5), Inches(3.0), Pt(4))
tb.fill.solid(); tb.fill.fore_color.rgb = SLATE; no_line(tb)
add_textbox(s, Inches(1.0), Inches(2.5), Inches(3.0), Inches(2.4),
            "OpenCLIP\nViT-B-32\n(frozen)", SZ_BODY, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
teal_arrow(s, Inches(4.25), Inches(3.5), Inches(1.4), Inches(0.4), AMBER)
# Probe card
pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(2.85), Inches(3.0), Inches(1.7))
pb.fill.solid(); pb.fill.fore_color.rgb = CARD; thin_border(pb, BORDER, Pt(1))
ptb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(2.85), Inches(3.0), Pt(4))
ptb.fill.solid(); ptb.fill.fore_color.rgb = AMBER; no_line(ptb)
add_textbox(s, Inches(5.9), Inches(2.85), Inches(3.0), Inches(1.7),
            "nn.Linear\n(512 → 20)", SZ_BODY, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
teal_arrow(s, Inches(9.15), Inches(3.5), Inches(1.1), Inches(0.4), AMBER)
# Output card
ob = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.45), Inches(2.85), Inches(1.95), Inches(1.7))
ob.fill.solid(); ob.fill.fore_color.rgb = CARD; thin_border(ob, BORDER, Pt(1))
otb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.45), Inches(2.85), Inches(1.95), Pt(4))
otb.fill.solid(); otb.fill.fore_color.rgb = TEAL; no_line(otb)
add_textbox(s, Inches(10.45), Inches(2.85), Inches(1.95), Inches(1.7), "20 brands",
            SZ_BODY, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_caption(s, "Trained on Kaggle car images; 70/15/15 stratified split; early stopping at patience 7.", top=Inches(5.4))
add_note(s, "We freeze CLIP's general car knowledge and learn a lightweight 512→20 probe on top. Strong brand discrimination, no overfitting.")

# ── SLIDE 10: Video Brand Overlay ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, AMBER)
add_title(s, "Brands drawn live on the video")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), AMBER)
add_subtitle(s, "Car-only crops · 80×80 min · trucks excluded by design", Inches(1.55))
# Real output frame from Stage 2b (YOLO + CLIP brand overlay)
img_path = "/home/dongyuan/Desktop/computer_vision/other_document/clip_yolo_result_check.png"
pic = s.shapes.add_picture(img_path, MARGIN, Inches(2.2), width=Inches(7.8))
# Add a subtle border around the image
thin_border(pic, BORDER, Pt(1))
# Flow pills — the pipeline steps that produced this frame
flow = ["Frame", "YOLO crop", "CLIP", "softmax", "label"]
for i, item in enumerate(flow):
    pill(s, Inches(9.05), Inches(2.25) + i * Inches(0.78), Inches(3.35), Inches(0.58),
         item, AMBER)
add_caption(s, "Real output: YOLO26s boxes + CLIP brand chips (BMW, Peugeot …). Trucks excluded by design — car-only probe.")
add_note(s, "This is a real frame from the Stage 2b pipeline. YOLO26s detects the cars (YOLO conf=0.2); CLIP crops each car, runs it through the frozen ViT-B-32 backbone and the learned 512→20 linear probe, and overlays the top-1 brand with confidence. We only run CLIP on car crops above 80×80 px. Trucks are intentionally skipped — the probe was trained on car-only images and would miscalibrate on truck crops. Brand confidence is currently not thresholded (a known limit).")

# ── SLIDE 11: Divider Explain ─────────────────────────────────────────────────
s = divider_slide(prs, blank, "Explain", 3, INDIGO)
add_note(s, "We can locate and name every car. Now: how do we tell the story to a human who doesn't read bounding boxes?")

# ── SECTION 3 SPINE COLOR = INDIGO (Explain) ────────────────────────────────────
# ── SLIDE 12: From Labels to Language ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, INDIGO)
add_title(s, "The ultimate interface is plain language.")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), INDIGO)
# Real VLM output frame — the "Explain" stage in action
img_vlm1 = "/home/dongyuan/Desktop/computer_vision/other_document/vlm_caption_red_car_rain.png"
pic1 = s.shapes.add_picture(img_vlm1, Inches(2.0), Inches(2.0), width=Inches(9.3))
thin_border(pic1, BORDER, Pt(1))
# Caption banner label below the image
add_textbox(s, Inches(2.0), Inches(6.15), Inches(9.3), Inches(0.5),
            '"A red car is driving down a wet road in the rain, approaching a traffic light."',
            Pt(14), INDIGO, align=PP_ALIGN.CENTER)
add_caption(s, "Scene captioning + frame-based Q&A so any viewer — not just an engineer — can follow the story.")
add_note(s, "This is a real output from Stage 3. YOLO sees 'car, trafficLight' — the VLM narrates 'A red car is driving down a wet road in the rain, approaching a traffic light.' Boxes and brands are intermediate outputs. The real product is a sentence any human can read. That's what a customer demo needs.")

# ── SLIDE 13: VLM Adaptive Captioning ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, INDIGO)
add_title(s, "Adaptive captioning — only when the scene changes")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), INDIGO)
add_subtitle(s, "Qwen3-VL:4b via Ollama · adaptive gating · diff ≥ 12 · min interval 5 s", Inches(1.55))
# Two real VLM outputs side-by-side — demonstrating adaptive gating (different scenes → different captions)
img_vlm_a = "/home/dongyuan/Desktop/computer_vision/other_document/vlm_caption_red_car_rain.png"
img_vlm_b = "/home/dongyuan/Desktop/computer_vision/other_document/vlm_caption_truck_pedestrian.png"

# Left frame
x1 = MARGIN
pic_a = s.shapes.add_picture(img_vlm_a, x1, Inches(2.2), width=Inches(5.6))
thin_border(pic_a, BORDER, Pt(1))
add_textbox(s, x1, Inches(6.15), Inches(5.6), Inches(0.5),
            '"A red car is driving down a wet road in the rain, approaching a traffic light."',
            Pt(12), INDIGO, align=PP_ALIGN.CENTER)

# Right frame
x2 = MARGIN + Inches(5.6) + Inches(0.33)
pic_b = s.shapes.add_picture(img_vlm_b, x2, Inches(2.2), width=Inches(5.6))
thin_border(pic_b, BORDER, Pt(1))
add_textbox(s, x2, Inches(6.15), Inches(5.6), Inches(0.5),
            '"A white truck and a pedestrian are crossing the road at a pedestrian crossing..."',
            Pt(12), INDIGO, align=PP_ALIGN.CENTER)

add_caption(s, "Captions fire only when the scene actually changes — no redundant narration.")
add_note(s, "Two real frames from Stage 3, showing adaptive captioning in action. The left frame triggers a 'rainy driving' caption; the right frame triggers a 'pedestrian crossing' caption — because the scene changed enough (diff ≥ 12) and the minimum 5-second interval elapsed. The VLM produces fluent sentences with colour, weather, action, and spatial relationships — not just keyword lists. This is the difference between a detector output and a human-readable story.")

# ── SLIDE 14: Q&A Interface ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
section_spine(s, INDIGO)
add_title(s, "Frame-based Q&A — ask what you see")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), INDIGO)
video_placeholder(s, Inches(1.6), Inches(1.9), Inches(10.1), Inches(4.2),
                  "VLM Q&A interface — user typing a question, answer appearing (5-sec clip)")
add_caption(s, "Useful for accident reconstruction, driver training, or live demo interaction.")
add_note(s, "Pause any frame and ask: 'Why did the model flag this?' Turns a black-box detector into an explainable system.")

# ── SLIDE 15: Pipeline Architecture ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title(s, "Four stages · one coherent story")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
add_subtitle(s, "The main pipeline runs end-to-end; SegFormer runs alongside, independently.", Inches(1.55))

# ── Main pipeline row ──
stages = [("Raw\nVideo", LIGHT_SLATE), ("YOLO26s\nDetection", TEAL),
          ("CLIP Brand\nProbe", AMBER), ("VLM Caption\n& Q&A", INDIGO),
          ("Structured\nUnderstanding", NAVY)]
st_w = Inches(2.05); st_h = Inches(1.5); st_gap = Inches(0.42)
row_y = Inches(2.45)
total = 5 * st_w + 4 * st_gap
sx = (PRS_WIDTH - total) / 2
card_centers = []
for i, (text, color) in enumerate(stages):
    x = sx + i * (st_w + st_gap)
    card_centers.append(x + st_w / 2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, row_y, st_w, st_h)
    box.fill.solid(); box.fill.fore_color.rgb = CARD; thin_border(box, BORDER, Pt(1))
    tbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, row_y, st_w, Pt(4))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = color; no_line(tbar)
    add_textbox(s, x, row_y + Inches(0.1), st_w, st_h - Inches(0.1), text, Pt(13), NAVY,
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(stages) - 1:
        ax = x + st_w + Inches(0.06)
        teal_arrow(s, ax, row_y + st_h / 2 - Inches(0.13), st_gap - Inches(0.12), Inches(0.26))

# ── Independent SegFormer lane (below, visually offset) ──
branch_y = Inches(4.75)
bx = sx
branch = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, branch_y, st_w, st_h)
branch.fill.solid(); branch.fill.fore_color.rgb = CARD; thin_border(branch, BORDER, Pt(1))
bbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, branch_y, st_w, Pt(4))
bbar.fill.solid(); bbar.fill.fore_color.rgb = TEAL; no_line(bbar)
add_textbox(s, bx, branch_y + Inches(0.1), st_w, st_h - Inches(0.1),
            "SegFormer\nSegmentation", Pt(13), NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Dotted-style connector cue: a small vertical link from Raw Video down to SegFormer
link = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_centers[0] - Inches(0.01),
                          row_y + st_h, Inches(0.025), branch_y - (row_y + st_h))
link.fill.solid(); link.fill.fore_color.rgb = BORDER; no_line(link)
# Annotation to the right of the branch
add_textbox(s, bx + st_w + Inches(0.4), branch_y + Inches(0.25), Inches(6.5), Inches(1.0),
            "Standalone quality demo — shares the input video,\nbut feeds no downstream stage.",
            SZ_BODY, SLATE, anchor=MSO_ANCHOR.MIDDLE)

add_caption(s, "Speed: YOLO is fast · CLIP is moderate (one crop at a time) · SegFormer & VLM are slow, selective, or offline-only.", top=Inches(6.65))
add_note(s, "This is the full stack. Detect gives us structure. Enrich adds identity. Explain produces the narrative. SegFormer shares the same input video but runs on its own — it's a quality demo, not part of the live pipeline. Three sections, four stages, one demo video.")

# ── SLIDE 16: Results & Key Metrics ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_title(s, "Results & Key Metrics")
accent_rule(s, MARGIN, Inches(1.35), Inches(1.4), TEAL)
metrics = [("Detection", "P = 0.874\nR = 0.754\nmAP@0.5 = 0.842\nmAP@0.5:0.95 = 0.515", TEAL),
           ("Enrichment", "Probe accuracy ≈ 80 %\n20 brands\n(79.8 % on test split)", AMBER),
           ("Segmentation", "19 Cityscapes classes\n84 % mIoU (pretrained)", TEAL),
           ("Narration", "Adaptive captioning\nlocal inference", INDIGO)]
m_w = Inches(2.7); m_gap = Inches(0.32)
total = 4 * m_w + 3 * m_gap
mx = (PRS_WIDTH - total) / 2
for i, (title, body, color) in enumerate(metrics):
    x = mx + i * (m_w + m_gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), m_w, Inches(2.3))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; thin_border(card, BORDER, Pt(1))
    tbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), m_w, Pt(5))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = color; no_line(tbar)
    add_textbox(s, x, Inches(2.45), m_w, Inches(0.5), title, SZ_SUBTITLE, color,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.15), Inches(3.05), m_w - Inches(0.3), Inches(1.2),
                body, Pt(13), SLATE, align=PP_ALIGN.CENTER)
add_caption(s, "YOLOv10n baseline → YOLO26s:  +24.7% mAP@0.5,  +32.7% mAP@0.5:0.95,  +26.1% recall.  Hardest class: pedestrian (mAP 0.692).", top=Inches(5.5))
add_note(s, "Detection: YOLOv10n (P=0.801, R=0.598, mAP@0.5=0.675, mAP@0.5:0.95=0.388) upgraded to YOLO26s (P=0.874, R=0.754, mAP@0.5=0.842, mAP@0.5:0.95=0.515) — a substantial accuracy lift across all metrics. The C2PSA backbone (9.4 M params) generalises well on 11 dashcam classes; pedestrian remains hardest due to scale and occlusion.  Enrichment: Brand recognition jumps from 55% zero-shot to 80% with the probe — a 25-point gain from a single linear layer on frozen embeddings. Segmentation leverages a pretrained SOTA model for immediate quality.")

# ── SLIDE 17: Future Work & Takeaway ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_slide_bg(s)
add_textbox(s, MARGIN, Inches(1.7), CONTENT_W, Inches(1.0),
            "From pixels to plain language —\nand next, to real-time edge deployment.",
            Pt(30), NAVY, bold=True, align=PP_ALIGN.CENTER)
accent_rule(s, Inches(5.97), Inches(3.0), Inches(1.4), TEAL)
items = ["Multi-camera fusion", "On-device quantization for edge speed", "Temporal tracking to reduce flicker"]
for i, text in enumerate(items):
    add_textbox(s, Inches(3.5), Inches(3.4) + i * Inches(0.62), Inches(6.5), Inches(0.5),
                f"·   {text}", SZ_BODY, SLATE)
add_note(s, "The pipeline works end-to-end on a local GPU. The next step is shrinking it: TensorRT, ONNX, and temporal smoothing so it runs on an embedded dashcam unit. Thank you — we're happy to take questions.")

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "/home/dongyuan/Desktop/computer_vision/Computer_Vision_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
